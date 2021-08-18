# -*- coding: utf-8 -*-
"""PowerPoint_Python

Automatically generated by Colaboratory.

Original file is located at
    https://colab.research.google.com/drive/1IHbdFy3mNeXExdYmiKZSKdDHjkg0HVux
"""

!pip install python-pptx

from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]

slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, World!"
subtitle.text = "Teste python Power Point"

slide1_register = prs.slide_layouts[1]

slide1 = prs.slides.add_slide(slide1_register)

title1 = slide1.shapes.title
title1.text = "Novo título teste"

slide2_register = prs.slide_layouts[2]
slide2 = prs.slides.add_slide(slide2_register)
title2 = slide2.shapes.title
title2.text = "Teste com imagem"

img = "/content/img_power_point.png"

from_left = Inches(3)
from_top = Inches(2)
adiciona_img = slide2.shapes.add_picture(img, from_left, from_top)

prs.save('teste_imagem.pptx')